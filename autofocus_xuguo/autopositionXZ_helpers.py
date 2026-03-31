#!/usr/bin/env python3
"""
Helper functions for X/Y auto-position.

This module keeps the future logic class clean by providing:
- square 2D mapping with scan-compatible JSON output
- map-to-map offset fitting with a quality threshold
"""

from __future__ import annotations

import json
import math
import os
import csv
import time
from dataclasses import asdict, dataclass
from datetime import datetime
from pathlib import Path
from tempfile import NamedTemporaryFile
from typing import Any, Callable

import matplotlib.pyplot as plt
import numpy as np
from scipy import ndimage
from scipy.optimize import curve_fit
from skimage.registration import phase_cross_correlation
from pptx import Presentation
from pptx.util import Inches


class OperationStoppedError(RuntimeError):
    """Raised when user requests stopping a long-running autofocus/autoposition operation."""


class _ScanInfoEncoder(json.JSONEncoder):
    def default(self, obj):
        if isinstance(obj, np.ndarray):
            return obj.tolist()
        if isinstance(obj, np.floating):
            value = float(obj)
            if math.isnan(value):
                return "NaN"
            return value
        if isinstance(obj, np.integer):
            return int(obj)
        if isinstance(obj, np.bool_):
            return bool(obj)
        return super().default(obj)


@dataclass
class OffsetFitResult:
    success: bool
    offset_x: float | None
    offset_y: float | None
    shift_x_pixels: float | None
    shift_y_pixels: float | None
    quality: float
    threshold: float
    error: float | None
    message: str


@dataclass
class GaussianPeakFitResult:
    success: bool
    peak_position: float | None
    peak_value: float | None
    baseline: float | None
    amplitude: float | None
    sigma: float | None
    r2: float | None
    message: str


def run_autoposition_square_mapping(
    hardware,
    *,
    center_x: float,
    center_y: float,
    points_per_line: int,
    span: float,
    save_path: str | os.PathLike[str],
    name: str = "autoposition_map",
    settle_time_s: float = 0.0,
    return_to_center: bool = True,
    comments: str = "",
    line_finished_callback: Callable[[int, int], None] | None = None,
    should_stop: Callable[[], bool] | None = None,
    process_events: Callable[[], None] | None = None,
) -> dict[str, Any]:
    """
    Acquire a square 2D map with X as the fast axis and Y as the slow axis.

    The saved JSON follows the same general top-level structure used by
    `core/scan.py`, so it can be loaded by the existing scan reader/plot code.
    """
    if points_per_line < 2:
        raise ValueError("points_per_line must be at least 2.")
    if span <= 0:
        raise ValueError("span must be positive.")

    x_target_device = _get_required_attr(hardware, "x_target_device")
    x_channel = _get_required_attr(hardware, "x_channel")
    y_target_device = _get_required_attr(hardware, "y_target_device")
    y_channel = _get_required_attr(hardware, "y_channel")
    ref_target_device = _get_required_attr(hardware, "reference_target_device")
    ref_channel = _get_required_attr(hardware, "reference_channel")

    x_values = np.linspace(center_x - span / 2.0, center_x + span / 2.0, points_per_line)
    y_values = np.linspace(center_y - span / 2.0, center_y + span / 2.0, points_per_line)

    image = np.full((points_per_line, points_per_line), np.nan, dtype=float)

    mapping_error: Exception | None = None
    try:
        for y_index, y_value in enumerate(y_values):
            _check_stop(should_stop)
            hardware.move_absoluteY(float(y_value))
            if settle_time_s > 0:
                time.sleep(settle_time_s)

            for x_index, x_value in enumerate(x_values):
                _check_stop(should_stop)
                hardware.move_absoluteX(float(x_value))
                if settle_time_s > 0:
                    time.sleep(settle_time_s)
                image[y_index, x_index] = float(hardware.read_reference_value())
                if process_events is not None:
                    process_events()

            if line_finished_callback is not None:
                line_finished_callback(y_index + 1, points_per_line)
            if process_events is not None:
                process_events()
    except Exception as exc:
        mapping_error = exc
        raise
    finally:
        if return_to_center:
            try:
                hardware.move_absoluteX(float(center_x))
                hardware.move_absoluteY(float(center_y))
            except Exception:
                if mapping_error is None:
                    raise

    info = _build_autoposition_scan_info(
        name=name,
        x_values=x_values,
        y_values=y_values,
        image=image,
        x_full_channel=f"{x_target_device}_{x_channel}",
        y_full_channel=f"{y_target_device}_{y_channel}",
        reference_full_channel=f"{ref_target_device}_{ref_channel}",
        center_x=float(center_x),
        center_y=float(center_y),
        span=float(span),
        points_per_line=int(points_per_line),
        comments=comments,
    )

    json_path = _save_autoposition_json(
        info=info,
        save_path=save_path,
        name=name,
        reference_full_channel=f"{ref_target_device}_{ref_channel}",
        span=float(span),
        points_per_line=int(points_per_line),
    )

    return {
        "info": info,
        "json_path": str(json_path),
        "x_values": x_values,
        "y_values": y_values,
        "image": image,
    }


def fit_offset(
    reference_map: np.ndarray,
    new_map: np.ndarray,
    *,
    x_values: np.ndarray | None = None,
    y_values: np.ndarray | None = None,
    span: float | None = None,
    points_per_line: int | None = None,
    upsample_factor: int = 20,
    quality_threshold: float = 0.6,
) -> OffsetFitResult:
    """
    Fit the X/Y compensation offset between two maps.

    Returned offsets are the compensation values to apply to the scan center:
    - `offset_x` should be added to the X center
    - `offset_y` should be added to the Y center

    If fitting quality is below `quality_threshold`, no offset is returned.
    """
    reference_map = np.asarray(reference_map, dtype=float)
    new_map = np.asarray(new_map, dtype=float)

    if reference_map.shape != new_map.shape:
        raise ValueError("reference_map and new_map must have the same shape.")
    if reference_map.ndim != 2:
        raise ValueError("reference_map and new_map must both be 2D arrays.")

    pixel_size_x, pixel_size_y = _resolve_pixel_sizes(
        shape=reference_map.shape,
        x_values=x_values,
        y_values=y_values,
        span=span,
        points_per_line=points_per_line,
    )

    reference_norm = _normalize_map(reference_map)
    new_norm = _normalize_map(new_map)

    try:
        shift_yx, error, _ = phase_cross_correlation(
            reference_norm,
            new_norm,
            upsample_factor=upsample_factor,
            disambiguate=True,
            normalization=None,
        )
    except TypeError:
        shift_yx, error, _ = phase_cross_correlation(
            reference_norm,
            new_norm,
            upsample_factor=upsample_factor,
        )

    candidate_shifts = [
        (float(shift_yx[0]), float(shift_yx[1])),
        (-float(shift_yx[0]), -float(shift_yx[1])),
    ]

    best_quality = float("-inf")
    best_shift_y_pixels = candidate_shifts[0][0]
    best_shift_x_pixels = candidate_shifts[0][1]
    for candidate_shift_y, candidate_shift_x in candidate_shifts:
        shifted_new = ndimage.shift(
            new_norm,
            shift=(candidate_shift_y, candidate_shift_x),
            order=1,
            mode="constant",
            cval=np.nan,
        )
        candidate_quality = _overlap_correlation(reference_norm, shifted_new)
        if np.isfinite(candidate_quality) and candidate_quality > best_quality:
            best_quality = float(candidate_quality)
            best_shift_y_pixels = float(candidate_shift_y)
            best_shift_x_pixels = float(candidate_shift_x)

    shift_y_pixels = best_shift_y_pixels
    shift_x_pixels = best_shift_x_pixels
    quality = best_quality

    offset_x = -shift_x_pixels * pixel_size_x
    offset_y = -shift_y_pixels * pixel_size_y

    if not np.isfinite(quality):
        return OffsetFitResult(
            success=False,
            offset_x=None,
            offset_y=None,
            shift_x_pixels=shift_x_pixels,
            shift_y_pixels=shift_y_pixels,
            quality=float("nan"),
            threshold=float(quality_threshold),
            error=None if error is None else float(error),
            message="Offset fit failed because fit quality could not be evaluated.",
        )

    if quality < quality_threshold:
        return OffsetFitResult(
            success=False,
            offset_x=None,
            offset_y=None,
            shift_x_pixels=shift_x_pixels,
            shift_y_pixels=shift_y_pixels,
            quality=float(quality),
            threshold=float(quality_threshold),
            error=None if error is None else float(error),
            message=(
                f"Offset fit rejected: quality {quality:.3f} is below "
                f"threshold {quality_threshold:.3f}."
            ),
        )

    return OffsetFitResult(
        success=True,
        offset_x=float(offset_x),
        offset_y=float(offset_y),
        shift_x_pixels=shift_x_pixels,
        shift_y_pixels=shift_y_pixels,
        quality=float(quality),
        threshold=float(quality_threshold),
        error=None if error is None else float(error),
        message="Offset fit succeeded.",
    )


def offset_fit_to_dict(result: OffsetFitResult) -> dict[str, Any]:
    return asdict(result)


def run_autofocus_z_profile(
    hardware,
    *,
    z_start_um: float,
    z_end_um: float,
    step_um: float,
    settle_time_s: float = 0.0,
    point_finished_callback: Callable[[int, int], None] | None = None,
    should_stop: Callable[[], bool] | None = None,
    process_events: Callable[[], None] | None = None,
) -> dict[str, np.ndarray]:
    """
    Run a 1D autofocus profile scan in Z.

    This helper uses absolute-height moves and reads one reference value at each point.
    Scan direction is preserved from `z_start_um` to `z_end_um`.
    """
    if step_um <= 0:
        raise ValueError("step_um must be positive.")

    start_um = float(z_start_um)
    end_um = float(z_end_um)
    if abs(end_um - start_um) <= 1e-12:
        z_positions = np.asarray([start_um], dtype=float)
    else:
        direction = 1.0 if end_um > start_um else -1.0
        signed_step = abs(float(step_um)) * direction
        n_steps = int(np.floor(abs(end_um - start_um) / abs(float(step_um)))) + 1
        z_positions = start_um + np.arange(n_steps, dtype=float) * signed_step
        if abs(float(z_positions[-1]) - end_um) > 1e-12:
            z_positions = np.append(z_positions, end_um)

    values = np.full_like(z_positions, np.nan, dtype=float)
    for index, position_um in enumerate(z_positions):
        _check_stop(should_stop)
        hardware.move_absolute_height(float(position_um))
        if settle_time_s > 0:
            time.sleep(float(settle_time_s))
        values[index] = float(hardware.read_reference_value())
        if point_finished_callback is not None:
            point_finished_callback(index + 1, len(z_positions))
        if process_events is not None:
            process_events()

    return {"z_positions": z_positions, "values": values}


def fit_gaussian_peak(
    z_positions: np.ndarray,
    values: np.ndarray,
) -> GaussianPeakFitResult:
    """Fit a Gaussian peak on 1D autofocus profile data."""
    z_positions = np.asarray(z_positions, dtype=float)
    values = np.asarray(values, dtype=float)
    mask = np.isfinite(z_positions) & np.isfinite(values)
    if np.count_nonzero(mask) < 5:
        return GaussianPeakFitResult(
            success=False,
            peak_position=None,
            peak_value=None,
            baseline=None,
            amplitude=None,
            sigma=None,
            r2=None,
            message="Need at least 5 finite points for Gaussian fitting.",
        )

    x = z_positions[mask]
    y = values[mask]
    if np.nanmax(y) <= np.nanmin(y):
        return GaussianPeakFitResult(
            success=False,
            peak_position=None,
            peak_value=None,
            baseline=None,
            amplitude=None,
            sigma=None,
            r2=None,
            message="Profile has zero dynamic range.",
        )

    x_min = float(np.min(x))
    x_max = float(np.max(x))
    span = max(x_max - x_min, 1e-6)
    baseline_guess = float(np.min(y))
    amplitude_guess = float(np.max(y) - baseline_guess)
    mu_guess = float(x[np.argmax(y)])
    sigma_guess = max(span / 6.0, 1e-6)

    def gaussian(data_x, baseline, amplitude, mu, sigma):
        return baseline + amplitude * np.exp(-0.5 * ((data_x - mu) / sigma) ** 2)

    try:
        params, _ = curve_fit(
            gaussian,
            x,
            y,
            p0=[baseline_guess, amplitude_guess, mu_guess, sigma_guess],
            bounds=(
                [-np.inf, 0.0, x_min - span, 1e-9],
                [np.inf, np.inf, x_max + span, span * 3.0],
            ),
            maxfev=20000,
        )
    except Exception as exc:
        return GaussianPeakFitResult(
            success=False,
            peak_position=None,
            peak_value=None,
            baseline=None,
            amplitude=None,
            sigma=None,
            r2=None,
            message=f"Gaussian fit failed: {exc}",
        )

    baseline, amplitude, mu, sigma = [float(item) for item in params]
    fitted = gaussian(x, baseline, amplitude, mu, sigma)
    ss_res = float(np.sum((y - fitted) ** 2))
    ss_tot = float(np.sum((y - np.mean(y)) ** 2))
    r2 = float("nan") if ss_tot == 0 else 1.0 - (ss_res / ss_tot)

    if not np.isfinite(mu) or not np.isfinite(amplitude) or amplitude <= 0:
        return GaussianPeakFitResult(
            success=False,
            peak_position=None,
            peak_value=None,
            baseline=baseline,
            amplitude=amplitude,
            sigma=sigma,
            r2=r2,
            message="Fitted Gaussian peak is invalid.",
        )

    return GaussianPeakFitResult(
        success=True,
        peak_position=float(mu),
        peak_value=float(baseline + amplitude),
        baseline=baseline,
        amplitude=amplitude,
        sigma=sigma,
        r2=r2,
        message="Gaussian fit succeeded.",
    )


def gaussian_fit_to_dict(result: GaussianPeakFitResult) -> dict[str, Any]:
    return asdict(result)


def export_history_to_csv(
    history_rows: list[dict[str, Any]],
    csv_path: str | os.PathLike[str],
) -> str:
    """Save a list of history dictionaries to CSV."""
    path = Path(csv_path)
    path.parent.mkdir(parents=True, exist_ok=True)

    if not history_rows:
        fieldnames = ["timestamp", "note"]
        rows = [{"timestamp": datetime.now().isoformat(timespec="seconds"), "note": "empty"}]
    else:
        fieldnames = sorted({key for row in history_rows for key in row.keys()})
        rows = history_rows

    with path.open("w", newline="", encoding="utf-8") as handle:
        writer = csv.DictWriter(handle, fieldnames=fieldnames)
        writer.writeheader()
        for row in rows:
            writer.writerow({name: row.get(name) for name in fieldnames})

    return str(path)


def append_autoposition_report_slide(
    *,
    ppt_path: str | os.PathLike[str],
    reference_image: np.ndarray,
    new_image: np.ndarray,
    offset_x: float | None,
    offset_y: float | None,
    quality: float,
    success: bool,
    details: str = "",
) -> str:
    """Append one autoposition comparison slide to a PowerPoint file."""
    status = "SUCCESS" if success else "FAILED"
    details = details.strip()

    fig, axes = plt.subplots(1, 2, figsize=(10.0, 4.5), constrained_layout=True)
    im0 = axes[0].imshow(np.asarray(reference_image, dtype=float), origin="lower", aspect="auto")
    axes[0].set_title("Reference map")
    axes[0].set_xlabel("X")
    axes[0].set_ylabel("Y")
    fig.colorbar(im0, ax=axes[0], fraction=0.046, pad=0.04)

    im1 = axes[1].imshow(np.asarray(new_image, dtype=float), origin="lower", aspect="auto")
    axes[1].set_title("Current map")
    axes[1].set_xlabel("X")
    axes[1].set_ylabel("Y")
    fig.colorbar(im1, ax=axes[1], fraction=0.046, pad=0.04)

    subtitle = (
        f"Offset X={_fmt(offset_x)}  Y={_fmt(offset_y)}  "
        f"Quality={quality:.4f}  Status={status}"
    )
    if details:
        subtitle = f"{subtitle}\n{details}"
    fig.suptitle("AutoPosition X/Y", fontsize=14)
    fig.text(0.5, 0.01, subtitle, ha="center", va="bottom", fontsize=10)

    return _append_figure_slide(
        ppt_path=ppt_path,
        title=f"AutoPosition X/Y ({status})",
        subtitle=datetime.now().isoformat(timespec="seconds"),
        figure=fig,
    )


def append_autofocus_report_slide(
    *,
    ppt_path: str | os.PathLike[str],
    coarse_positions: np.ndarray,
    coarse_values: np.ndarray,
    fine_positions: np.ndarray,
    fine_values: np.ndarray,
    coarse_peak: float | None,
    fine_peak: float | None,
    success: bool,
    details: str = "",
) -> str:
    """Append one autofocus profile slide to a PowerPoint file."""
    status = "SUCCESS" if success else "FAILED"
    details = details.strip()

    fig, axes = plt.subplots(1, 2, figsize=(10.0, 4.5), constrained_layout=True)
    axes[0].plot(np.asarray(coarse_positions, dtype=float), np.asarray(coarse_values, dtype=float), "o-", ms=3)
    if coarse_peak is not None:
        axes[0].axvline(float(coarse_peak), color="red", linestyle="--", label=f"peak={coarse_peak:.3f}")
        axes[0].legend(loc="best")
    axes[0].set_title("Coarse scan")
    axes[0].set_xlabel("Z (um)")
    axes[0].set_ylabel("Reference")
    axes[0].grid(True, alpha=0.3)

    axes[1].plot(np.asarray(fine_positions, dtype=float), np.asarray(fine_values, dtype=float), "o-", ms=3)
    if fine_peak is not None:
        axes[1].axvline(float(fine_peak), color="red", linestyle="--", label=f"peak={fine_peak:.3f}")
        axes[1].legend(loc="best")
    axes[1].set_title("Fine scan")
    axes[1].set_xlabel("Z (um)")
    axes[1].set_ylabel("Reference")
    axes[1].grid(True, alpha=0.3)

    subtitle = (
        f"Coarse peak={_fmt(coarse_peak)}  Fine peak={_fmt(fine_peak)}  Status={status}"
    )
    if details:
        subtitle = f"{subtitle}\n{details}"
    fig.suptitle("AutoFocus Z", fontsize=14)
    fig.text(0.5, 0.01, subtitle, ha="center", va="bottom", fontsize=10)

    return _append_figure_slide(
        ppt_path=ppt_path,
        title=f"AutoFocus Z ({status})",
        subtitle=datetime.now().isoformat(timespec="seconds"),
        figure=fig,
    )


def _build_autoposition_scan_info(
    *,
    name: str,
    x_values: np.ndarray,
    y_values: np.ndarray,
    image: np.ndarray,
    x_full_channel: str,
    y_full_channel: str,
    reference_full_channel: str,
    center_x: float,
    center_y: float,
    span: float,
    points_per_line: int,
    comments: str,
) -> dict[str, Any]:
    x_values = np.asarray(x_values, dtype=float)
    y_values = np.asarray(y_values, dtype=float)
    image = np.asarray(image, dtype=float)

    x_step = float(x_values[1] - x_values[0])
    y_step = float(y_values[1] - y_values[0])

    level0_data = np.full((1, points_per_line, points_per_line), np.nan, dtype=float)
    level0_data[0, :, :] = image
    level1_data = np.full((1, points_per_line), np.nan, dtype=float)

    timestamp = datetime.now().isoformat(timespec="seconds")

    return {
        "name": name,
        "levels": {
            "level0": {
                "setters": {
                    "setter0": {
                        "channel": x_full_channel,
                        "explicit": False,
                        "linear_setting": {
                            "start": float(x_values[0]),
                            "end": float(x_values[-1]),
                            "step": x_step,
                            "mid": float(center_x),
                            "span": float(span),
                            "points": int(points_per_line),
                            "destinations": x_values.copy(),
                        },
                        "explicit_setting": [],
                        "destinations": x_values.copy(),
                    }
                },
                "setting_method": "A",
                "getters": [reference_full_channel],
                "setting_array": np.array([x_values], dtype=float),
                "manual_set_before": [],
                "manual_set_after": [],
            },
            "level1": {
                "setters": {
                    "setter0": {
                        "channel": y_full_channel,
                        "explicit": False,
                        "linear_setting": {
                            "start": float(y_values[0]),
                            "end": float(y_values[-1]),
                            "step": y_step,
                            "mid": float(center_y),
                            "span": float(span),
                            "points": int(points_per_line),
                            "destinations": y_values.copy(),
                        },
                        "explicit_setting": [],
                        "destinations": y_values.copy(),
                    }
                },
                "setting_method": "A",
                "getters": ["none"],
                "setting_array": np.array([y_values], dtype=float),
                "manual_set_before": [],
                "manual_set_after": [],
            },
        },
        "data": [level0_data, level1_data],
        "plots": {
            "line_plots": {},
            "image_plots": {
                "0": {
                    "x": "level0",
                    "y": "level1",
                    "z": f"L0G0_{reference_full_channel}",
                }
            },
        },
        "plots_per_page": "2x1",
        "autoposition_metadata": {
            "saved_at": timestamp,
            "center_x": float(center_x),
            "center_y": float(center_y),
            "span": float(span),
            "points_per_line": int(points_per_line),
            "x_channel": x_full_channel,
            "y_channel": y_full_channel,
            "reference_channel": reference_full_channel,
            "comments": comments,
        },
    }


def _save_autoposition_json(
    *,
    info: dict[str, Any],
    save_path: str | os.PathLike[str],
    name: str,
    reference_full_channel: str,
    span: float,
    points_per_line: int,
) -> Path:
    autoposition_dir = Path(save_path) / "autoposition"
    autoposition_dir.mkdir(parents=True, exist_ok=True)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    base_name = (
        f"{timestamp}_{_sanitize(name)}"
        f"_ref-{_sanitize(reference_full_channel)}"
        f"_span-{_sanitize(f'{span:.6g}')}"
        f"_pts-{points_per_line}.json"
    )
    json_path = autoposition_dir / base_name

    with json_path.open("w", encoding="utf-8") as handle:
        json.dump(info, handle, cls=_ScanInfoEncoder, indent=4)

    return json_path


def _append_figure_slide(
    *,
    ppt_path: str | os.PathLike[str],
    title: str,
    subtitle: str,
    figure,
) -> str:
    ppt_file = Path(ppt_path)
    ppt_file.parent.mkdir(parents=True, exist_ok=True)

    with NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        temp_path = Path(tmp.name)
    try:
        figure.savefig(temp_path, dpi=160)
        if _try_append_figure_slide_via_com(
            ppt_path=str(ppt_file),
            title=str(title),
            subtitle=str(subtitle),
            image_path=str(temp_path),
        ):
            return str(ppt_file)

        if ppt_file.exists():
            presentation = Presentation(str(ppt_file))
        else:
            presentation = Presentation()
            presentation.slide_width = Inches(13.333)
            presentation.slide_height = Inches(7.5)

        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.35))
        title_box.text_frame.text = str(title)
        subtitle_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.52), Inches(12.7), Inches(0.28))
        subtitle_box.text_frame.text = str(subtitle)

        available_w = float(slide_width) - float(Inches(0.6))
        available_h = float(slide_height) - float(Inches(1.1))
        left = float(Inches(0.3))
        top = float(Inches(0.9))
        slide.shapes.add_picture(
            str(temp_path),
            int(left),
            int(top),
            width=int(available_w),
            height=int(available_h),
        )
        presentation.save(str(ppt_file))
    finally:
        plt.close(figure)
        if temp_path.exists():
            try:
                temp_path.unlink()
            except OSError:
                pass

    return str(ppt_file)


def _resolve_pixel_sizes(
    *,
    shape: tuple[int, int],
    x_values: np.ndarray | None,
    y_values: np.ndarray | None,
    span: float | None,
    points_per_line: int | None,
) -> tuple[float, float]:
    height, width = shape

    if x_values is not None:
        x_values = np.asarray(x_values, dtype=float)
        if x_values.ndim != 1 or x_values.size != width:
            raise ValueError("x_values must be a 1D array matching map width.")
        pixel_size_x = float(x_values[1] - x_values[0])
    else:
        if span is None or points_per_line is None:
            raise ValueError(
                "Provide x_values/y_values or provide span and points_per_line."
            )
        if int(points_per_line) != width or int(points_per_line) != height:
            raise ValueError(
                "points_per_line must match the map shape when x_values/y_values "
                "are not provided."
            )
        pixel_size_x = float(span) / float(points_per_line - 1)

    if y_values is not None:
        y_values = np.asarray(y_values, dtype=float)
        if y_values.ndim != 1 or y_values.size != height:
            raise ValueError("y_values must be a 1D array matching map height.")
        pixel_size_y = float(y_values[1] - y_values[0])
    else:
        if span is None or points_per_line is None:
            raise ValueError(
                "Provide x_values/y_values or provide span and points_per_line."
            )
        if int(points_per_line) != width or int(points_per_line) != height:
            raise ValueError(
                "points_per_line must match the map shape when x_values/y_values "
                "are not provided."
            )
        pixel_size_y = float(span) / float(points_per_line - 1)

    return pixel_size_x, pixel_size_y


def _normalize_map(image: np.ndarray) -> np.ndarray:
    finite = np.isfinite(image)
    if finite.sum() < 4:
        raise ValueError("Map must contain at least 4 finite values.")

    values = image[finite]
    std = float(np.std(values))
    if std == 0:
        raise ValueError("Map has zero variance and cannot be fit.")

    mean = float(np.mean(values))
    normalized = np.full_like(image, np.nan, dtype=float)
    normalized[finite] = (values - mean) / std
    return normalized


def _overlap_correlation(reference: np.ndarray, shifted: np.ndarray) -> float:
    mask = np.isfinite(reference) & np.isfinite(shifted)
    if mask.sum() < 4:
        return float("nan")

    ref_values = reference[mask]
    shifted_values = shifted[mask]
    corr = np.corrcoef(ref_values, shifted_values)[0, 1]
    return float(corr)


def _get_required_attr(obj: Any, name: str) -> str:
    value = getattr(obj, name, None)
    if value is None or value == "":
        raise RuntimeError(f"Hardware attribute '{name}' is not configured.")
    return str(value)


def _sanitize(text: str) -> str:
    safe = []
    for char in str(text):
        if char.isalnum() or char in ("-", "_"):
            safe.append(char)
        elif char in (".",):
            safe.append("p")
        else:
            safe.append("-")
    return "".join(safe).strip("-") or "item"


def _fmt(value: float | None) -> str:
    if value is None or not np.isfinite(value):
        return "nan"
    return f"{float(value):.6g}"


def _check_stop(should_stop: Callable[[], bool] | None) -> None:
    if should_stop is not None and bool(should_stop()):
        raise OperationStoppedError("Operation stopped by user request.")


def _try_append_figure_slide_via_com(
    *,
    ppt_path: str,
    title: str,
    subtitle: str,
    image_path: str,
) -> bool:
    if os.name != "nt":
        return False
    try:
        import win32com.client  # type: ignore
    except Exception:
        return False

    ppt_app = None
    presentation = None
    try:
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = True
        full_path = os.path.abspath(ppt_path)

        if not os.path.exists(full_path):
            presentation = ppt_app.Presentations.Add()
            presentation.SaveAs(full_path)
        else:
            for pres in ppt_app.Presentations:
                if os.path.normcase(pres.FullName) == os.path.normcase(full_path):
                    presentation = pres
                    break
            if presentation is None:
                presentation = ppt_app.Presentations.Open(
                    full_path, ReadOnly=False, Untitled=False, WithWindow=True
                )

        slide = presentation.Slides.Add(presentation.Slides.Count + 1, 12)
        slide_width = float(presentation.PageSetup.SlideWidth)
        slide_height = float(presentation.PageSetup.SlideHeight)

        title_box = slide.Shapes.AddTextbox(1, 20, 10, slide_width - 40, 24)
        title_box.TextFrame.TextRange.Text = title
        title_box.TextFrame.TextRange.Font.Bold = True
        title_box.TextFrame.TextRange.Font.Size = 16

        subtitle_box = slide.Shapes.AddTextbox(1, 20, 36, slide_width - 40, 18)
        subtitle_box.TextFrame.TextRange.Text = subtitle
        subtitle_box.TextFrame.TextRange.Font.Size = 11

        margin = 20.0
        image_left = margin
        image_top = 60.0
        image_width = slide_width - (2 * margin)
        image_height = slide_height - image_top - margin
        slide.Shapes.AddPicture(
            FileName=os.path.abspath(image_path),
            LinkToFile=False,
            SaveWithDocument=True,
            Left=image_left,
            Top=image_top,
            Width=image_width,
            Height=image_height,
        )
        presentation.Save()
        return True
    except Exception:
        return False
    finally:
        if presentation is not None:
            del presentation
        if ppt_app is not None:
            del ppt_app
